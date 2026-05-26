"""
UNSW-NB15 Saldırı Tespiti - PowerPoint Sunum Oluşturucu
python-pptx ile 14 slaytlık akademik sunum
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

BASE = r"C:\Users\seyma\OneDrive\Masaüstü\bigdataproje"
OUT  = os.path.join(BASE, "BigVeri_SaldiriTespiti_Sunum.pptx")

# ─── RENKLER ─────────────────────────────────────────────────
C_DARK_BLUE   = RGBColor(0x1A, 0x23, 0x7E)
C_MED_BLUE    = RGBColor(0x28, 0x35, 0x93)
C_ACCENT_BLUE = RGBColor(0x15, 0x65, 0xC0)
C_LIGHT_BLUE  = RGBColor(0xE3, 0xF2, 0xFD)
C_SKY_BLUE    = RGBColor(0xBB, 0xDE, 0xFB)
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK_TEXT   = RGBColor(0x1A, 0x23, 0x7E)
C_GRAY_TEXT   = RGBColor(0x55, 0x65, 0x81)
C_GREEN       = RGBColor(0x2E, 0x7D, 0x32)
C_GREEN_LIGHT = RGBColor(0xE8, 0xF5, 0xE9)
C_RED         = RGBColor(0xC6, 0x28, 0x28)
C_RED_LIGHT   = RGBColor(0xFF, 0xEB, 0xEE)
C_ORANGE      = RGBColor(0xE6, 0x5C, 0x00)
C_ORANGE_LIGHT= RGBColor(0xFF, 0xF3, 0xE0)
C_PURPLE      = RGBColor(0x4A, 0x14, 0x8C)
C_PURPLE_LIGHT= RGBColor(0xF3, 0xE5, 0xF5)

SW = 13.333   # slide width  (inches)
SH = 7.5      # slide height (inches)
HEADER_H = 0.85

def rgb_hex(rgb: RGBColor) -> str:
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"

# ─── YARDIMCI FONKSİYONLAR ───────────────────────────────────
def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = Inches(SW)
    prs.slide_height = Inches(SH)
    return prs

def blank_slide(prs):
    blank = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank)

def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_w=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, text, x, y, w, h,
                 font_size=14, bold=False, italic=False,
                 color=None, align=PP_ALIGN.LEFT,
                 font_face="Calibri", wrap=True,
                 valign=None, margin_inches=0.05):
    if color is None:
        color = C_DARK_TEXT
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if valign:
        tf.vertical_anchor = valign
    # set internal margin
    txBox.text_frame.margin_left   = Inches(margin_inches)
    txBox.text_frame.margin_right  = Inches(margin_inches)
    txBox.text_frame.margin_top    = Inches(margin_inches * 0.5)
    txBox.text_frame.margin_bottom = Inches(margin_inches * 0.5)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_face
    return txBox

def add_rich_text(slide, items, x, y, w, h,
                  wrap=True, valign=None, margin_inches=0.05):
    """items: list of (text, dict_of_options, is_new_para)"""
    from pptx.enum.text import PP_ALIGN
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.margin_left   = Inches(margin_inches)
    tf.margin_right  = Inches(margin_inches)
    tf.margin_top    = Inches(margin_inches * 0.5)
    tf.margin_bottom = Inches(margin_inches * 0.5)
    first = True
    cur_para = None
    for (text, opts, new_para) in items:
        if new_para or first:
            if first:
                cur_para = tf.paragraphs[0]
                first = False
            else:
                cur_para = tf.add_paragraph()
            cur_para.alignment = opts.get('align', PP_ALIGN.LEFT)
            if 'space_before' in opts:
                cur_para.space_before = Pt(opts['space_before'])
            if 'space_after' in opts:
                cur_para.space_after = Pt(opts['space_after'])
        run = cur_para.add_run()
        run.text = text
        run.font.size    = Pt(opts.get('size', 14))
        run.font.bold    = opts.get('bold', False)
        run.font.italic  = opts.get('italic', False)
        run.font.name    = opts.get('face', 'Calibri')
        run.font.color.rgb = opts.get('color', C_DARK_TEXT)
    return txBox

def add_header(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = C_DARK_BLUE
    rect = add_rect(slide, 0, 0, SW, HEADER_H, bg_color)
    tb = add_text_box(slide, title_text,
                      0.3, 0.08, SW - 0.6, HEADER_H - 0.16,
                      font_size=26, bold=True, color=C_WHITE,
                      align=PP_ALIGN.LEFT, font_face="Calibri",
                      margin_inches=0)
    return rect, tb

def add_slide_number(slide, num):
    add_text_box(slide, str(num),
                 SW - 0.55, SH - 0.38, 0.45, 0.3,
                 font_size=10, color=C_GRAY_TEXT,
                 align=PP_ALIGN.RIGHT, margin_inches=0)

def set_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text

def add_image(slide, img_name, x, y, w, h):
    path = os.path.join(BASE, img_name)
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(x), Inches(y),
                                  Inches(w), Inches(h))
    else:
        print(f"  UYARI: {img_name} bulunamadı")

def add_bullet_box(slide, bullets, x, y, w, h,
                   font_size=14, color=None, bullet_color=None,
                   indent=0.25, line_space=6):
    if color is None:
        color = C_DARK_TEXT
    if bullet_color is None:
        bullet_color = C_ACCENT_BLUE
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left   = Inches(0.05)
    tf.margin_right  = Inches(0.05)
    tf.margin_top    = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(line_space)
        # manual bullet using a unicode dash for cleaner look
        run = p.add_run()
        run.text = f"▸  {b}"
        run.font.size  = Pt(font_size)
        run.font.name  = "Calibri"
        run.font.color.rgb = color
    return txBox

def metric_card(slide, x, y, w, h, label, value, bg_color, text_color=None):
    if text_color is None:
        text_color = C_WHITE
    add_rect(slide, x, y, w, h, bg_color)
    # label
    add_text_box(slide, label,
                 x + 0.08, y + 0.12, w - 0.16, h * 0.38,
                 font_size=12, bold=False, color=text_color,
                 align=PP_ALIGN.CENTER, margin_inches=0)
    # value
    add_text_box(slide, value,
                 x + 0.08, y + h * 0.42, w - 0.16, h * 0.50,
                 font_size=24, bold=True, color=text_color,
                 align=PP_ALIGN.CENTER, margin_inches=0)

def add_table_styled(slide, headers, rows, x, y, w, h,
                     header_bg=None, row_alt=None):
    if header_bg is None:
        header_bg = C_DARK_BLUE
    if row_alt is None:
        row_alt = C_LIGHT_BLUE
    all_rows = [headers] + rows
    n_rows = len(all_rows)
    n_cols = len(headers)
    tbl = slide.shapes.add_table(
        n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl.table.first_row = True
    col_w = w / n_cols
    for ci in range(n_cols):
        tbl.table.columns[ci].width = Inches(col_w)
    for ri, row_data in enumerate(all_rows):
        for ci, cell_text in enumerate(row_data):
            cell = tbl.table.cell(ri, ci)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            if ri == 0:
                run.font.bold  = True
                run.font.color.rgb = C_WHITE
                run.font.size  = Pt(12)
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            else:
                run.font.color.rgb = C_DARK_TEXT
                run.font.size  = Pt(11)
                cell.fill.solid()
                if ri % 2 == 0:
                    cell.fill.fore_color.rgb = row_alt
                else:
                    cell.fill.fore_color.rgb = C_WHITE
    return tbl

# ════════════════════════════════════════════════════════════
# SLAYT ÜRETİMİ
# ════════════════════════════════════════════════════════════
prs = new_prs()

# ── SLAYT 1: KAPAK ──────────────────────────────────────────
print("Slayt 1: Kapak")
sl = blank_slide(prs)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = C_DARK_BLUE

# Dekoratif arka plan şekilleri
add_rect(sl, 0, 0, SW * 0.6, SH, RGBColor(0x1E, 0x2E, 0x8E), None, 0)
add_rect(sl, 0, SH - 1.2, SW, 1.2, RGBColor(0x0D, 0x47, 0xA1), None, 0)
add_rect(sl, 0, SH - 1.25, SW, 0.05, C_ACCENT_BLUE, None, 0)

# Başlık
add_text_box(sl, "Büyük Veri Güvenliği:",
             0.6, 1.0, 10.5, 1.0,
             font_size=34, bold=True, color=C_WHITE,
             align=PP_ALIGN.LEFT, margin_inches=0)
add_text_box(sl, "Ağ Trafiği Verilerinde Makine Öğrenimi ile\nSaldırı Tespiti ve Etik Değerlendirme",
             0.6, 1.85, 10.5, 1.6,
             font_size=26, bold=False, color=RGBColor(0xBB, 0xDE, 0xFB),
             align=PP_ALIGN.LEFT, margin_inches=0)

# Alt çizgi
add_rect(sl, 0.6, 3.55, 7.0, 0.05, C_ACCENT_BLUE, None, 0)

add_text_box(sl, "UNSW-NB15 Veri Seti Üzerinde Random Forest Uygulaması",
             0.6, 3.7, 10.5, 0.6,
             font_size=17, bold=False, color=RGBColor(0x90, 0xCA, 0xF9),
             align=PP_ALIGN.LEFT, margin_inches=0)
add_text_box(sl, "Büyük Veri Dersi  |  Yüksek Lisans",
             0.6, SH - 1.1, 10.5, 0.55,
             font_size=14, bold=False, color=RGBColor(0xBB, 0xDE, 0xFB),
             align=PP_ALIGN.LEFT, margin_inches=0)

add_slide_number(sl, 1)
set_notes(sl, "Merhaba, bugün büyük veri güvenliği konusunda hazırladığım projeyi sunacağım. Bu çalışmada ağ trafiği verilerinde makine öğrenimi kullanarak saldırı tespiti yapıyor ve bunun etik boyutlarını değerlendiriyorum.")

# ── SLAYT 2: PROJENİN AMACI ─────────────────────────────────
print("Slayt 2: Projenin Amacı")
sl = blank_slide(prs)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = C_WHITE
add_header(sl, "Projenin Amacı")
add_slide_number(sl, 2)

bullets = [
    "Ağ trafiği verileri üzerinde makine öğrenimi ile saldırı tespiti modeli kurmak",
    "Modelin başarısını teknik metriklerle (Accuracy, F1, ROC-AUC) ölçmek",
    "Yanlış sınıflandırmaların (FP/FN) güvenlik ve etik risklerini analiz etmek",
    "İşlenen verilerin mahremiyet ve GDPR boyutunu değerlendirmek",
]
add_bullet_box(sl, bullets, 0.5, HEADER_H + 0.2, SW - 1.0, 3.2, font_size=17)

# Öne çıkan kutu
add_rect(sl, 0.5, SH - 1.9, SW - 1.0, 1.25, C_LIGHT_BLUE, C_ACCENT_BLUE, 1)
add_text_box(sl, "Neden Büyük Veri?",
             0.65, SH - 1.85, SW - 1.3, 0.38,
             font_size=14, bold=True, color=C_ACCENT_BLUE,
             align=PP_ALIGN.LEFT, margin_inches=0)
add_text_box(sl, "Ağ trafiği verileri saniyede milyonlarca kayıt üretir — yüksek hacim, hız ve çeşitlilik içerir. Bu ölçekte manuel analiz yetersiz kalır; makine öğrenimi zorunlu hale gelir.",
             0.65, SH - 1.50, SW - 1.3, 0.85,
             font_size=13, bold=False, color=C_DARK_TEXT,
             align=PP_ALIGN.LEFT, margin_inches=0)
set_notes(sl, "Bu projenin dört temel amacı var. Sadece teknik bir model kurmakla kalmıyoruz, aynı zamanda bu modelin yanlış kararlarının ne gibi sonuçlar doğurduğunu güvenlik ve etik açısından inceliyoruz. Ağ trafiği verileri saniyede milyonlarca kayıt ürettiği için manuel analiz mümkün değil — bu da büyük veri yöntemlerini zorunlu kılıyor.")

# ── SLAYT 3: VERİ SETİ ──────────────────────────────────────
print("Slayt 3: UNSW-NB15 Veri Seti")
sl = blank_slide(prs)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = C_WHITE
add_header(sl, "UNSW-NB15 Veri Seti")
add_slide_number(sl, 3)

# Sol sütun - bilgi
add_text_box(sl, "Veri Seti Hakkında",
             0.4, HEADER_H + 0.15, 5.5, 0.4,
             font_size=15, bold=True, color=C_ACCENT_BLUE, margin_inches=0)
bullets_l = [
    "Australian Centre for Cyber Security (ACCS) tarafından oluşturuldu",
    "IXIA PerfectStorm aracıyla gerçek saldırı simülasyonu",
    "100 GB ham trafik verisi (PCAP formatı)",
    "Toplam 2.5 milyon+ kayıt, 49 özellik",
    "9 farklı saldırı kategorisi barındırmaktadır",
    "Akademik alanda yaygın referans veri seti",
]
add_bullet_box(sl, bullets_l, 0.4, HEADER_H + 0.6, 5.5, 4.0, font_size=13)

# Sağ sütun - tablo
add_text_box(sl, "Saldırı Kategorileri",
             6.2, HEADER_H + 0.15, 6.7, 0.4,
             font_size=15, bold=True, color=C_ACCENT_BLUE, margin_inches=0)
headers = ["Saldırı Tipi", "Açıklama"]
rows = [
    ["Fuzzers",        "Sistem çöküşü için rastgele veri"],
    ["DoS",            "Hizmet aksatma saldırısı"],
    ["Backdoor",       "Yetkisiz erişim arka kapısı"],
    ["Exploits",       "Güvenlik açığı istismarı"],
    ["Reconnaissance", "Keşif ve tarama saldırısı"],
    ["Shellcode",      "Kabuk kodu enjeksiyonu"],
    ["Generic",        "Şifreleme kırma saldırısı"],
    ["Analysis",       "Port tarama ve analiz"],
    ["Worms",          "Ağ solucanı yayılımı"],
]
add_table_styled(sl, headers, rows, 6.2, HEADER_H + 0.6, 6.7, 5.2)
set_notes(sl, "UNSW-NB15, Avustralya Siber Güvenlik Merkezi tarafından gerçek dünya saldırıları simüle edilerek oluşturulmuş, akademik alanda yaygın kullanılan bir veri setidir. Toplam 2.5 milyonu aşkın kayıt içeriyor ve 9 farklı saldırı türünü barındırıyor. Bu veri setinin büyüklüğü ve çeşitliliği onu büyük veri çalışmalarına ideal yapıyor.")

# ── SLAYT 4: NEDEN BÜYÜK VERİ? (3V) ────────────────────────
print("Slayt 4: Neden Büyük Veri?")
sl = blank_slide(prs)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = C_WHITE
add_header(sl, "Neden Büyük Veri? — 3V Modeli")
add_slide_number(sl, 4)

cards = [
    ("HACİM\n(Volume)", "Milyonlarca ağ bağlantısı kaydı\nGerçek dünya ağ ortamı\n100 GB+ ham veri\nCSV + PCAP formatı", C_DARK_BLUE),
    ("HIZ\n(Velocity)",  "Gerçek zamanlı trafik akışı\nAnlık tehdit tespiti gerekliliği\nSaniyede binlerce kayıt\nSürekli güncellenen veri", C_ACCENT_BLUE),
    ("ÇEŞİTLİLİK\n(Variety)", "49 farklı ağ özelliği\n9 saldırı kategorisi\nProtokoller: TCP, UDP, HTTP...\nZamanlama + istatistiksel özellikler", C_MED_BLUE),
]
card_w = (SW - 1.2) / 3
for i, (title, body, clr) in enumerate(cards):
    cx = 0.4 + i * (card_w + 0.2)
    add_rect(sl, cx, HEADER_H + 0.15, card_w, 1.0, clr)
    add_text_box(sl, title,
                 cx + 0.1, HEADER_H + 0.18, card_w - 0.2, 0.9,
                 font_size=16, bold=True, color=C_WHITE,
                 align=PP_ALIGN.CENTER, margin_inches=0)
    add_rect(sl, cx, HEADER_H + 1.18, card_w, 3.2, C_LIGHT_BLUE, clr, 1)
    add_text_box(sl, body,
                 cx + 0.15, HEADER_H + 1.25, card_w - 0.3, 3.05,
                 font_size=13, color=C_DARK_TEXT,
                 align=PP_ALIGN.LEFT, margin_inches=0)

# Alt banner
add_rect(sl, 0.4, SH - 1.05, SW - 0.8, 0.7, RGBColor(0x0D, 0x47, 0xA1))
add_text_box(sl, "→  Manuel analiz bu ölçekte imkânsızdır. Makine öğrenimi zorunludur.",
             0.6, SH - 1.00, SW - 1.2, 0.58,
             font_size=15, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER, margin_inches=0)
set_notes(sl, "Büyük verinin 3V modelini bu veri setiyle örnekleyebiliriz. Hacim olarak milyonlarca kayıt, hız olarak gerçek zamanlı trafik akışı, çeşitlilik olarak ise 49 farklı özellik ve 9 saldırı türü. Bu ölçekte ve hızda insan analizi mümkün değil — bu yüzden modellemek zorunlu hale geliyor.")

# ── SLAYT 5: EDA - SINIF DAĞILIMI ───────────────────────────
print("Slayt 5: EDA - Sınıf Dağılımı")
sl = blank_slide(prs)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = C_WHITE
add_header(sl, "Veri Keşfi (EDA): Sınıf Dağılımı")
add_slide_number(sl, 5)

add_image(sl, "class_distribution.png",
          0.4, HEADER_H + 0.1, SW - 0.8, 4.6)

# Alt istatistik kutuları
stat_w = (SW - 1.0) / 3 - 0.1
stats = [
    ("Eğitim Seti",    "82,332 kayıt",   C_ACCENT_BLUE),
    ("Test Seti",      "175,341 kayıt",  C_MED_BLUE),
    ("Saldırı Türleri","9 kategori",     C_DARK_BLUE),
]
for i, (lbl, val, clr) in enumerate(stats):
    sx = 0.5 + i * (stat_w + 0.15)
    add_rect(sl, sx, SH - 0.95, stat_w, 0.7, clr)
    add_text_box(sl, lbl, sx + 0.05, SH - 0.93, stat_w - 0.1, 0.28,
                 font_size=10, color=C_WHITE, align=PP_ALIGN.CENTER, margin_inches=0)
    add_text_box(sl, val, sx + 0.05, SH - 0.67, stat_w - 0.1, 0.35,
                 font_size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, margin_inches=0)
set_notes(sl, "Eğitim setimizde kayıtların yaklaşık yüzde otuz yedisi normal, geri kalanı saldırı trafiğidir. Saldırı kategorileri arasında Fuzzers, Reconnaissance ve Generic en yaygın olanlardır. Bu sınıf dengesizliği, modelin class_weight parametresiyle dengelenmesini gerektirdi.")

# ── SLAYT 6: EDA - PROTOKOL & KORELASYON ────────────────────
print("Slayt 6: EDA - Protokol & Korelasyon")
sl = blank_slide(prs)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = C_WHITE
add_header(sl, "Veri Keşfi: Protokol Dağılımı ve Korelasyon")
add_slide_number(sl, 6)

add_image(sl, "protocol_state.png",   0.3, HEADER_H + 0.05, SW - 0.6, 3.1)
add_image(sl, "correlation_heatmap.png", 0.3, HEADER_H + 3.25, SW - 0.6, 2.95)

add_text_box(sl, "▸  UDP ve TCP protokolleri dominant  |  ct_state_ttl, sttl, rate özellikleri hedefle en güçlü korelasyona sahip",
             0.5, SH - 0.38, SW - 1.0, 0.32,
             font_size=11, color=C_GRAY_TEXT,
             align=PP_ALIGN.CENTER, margin_inches=0)
set_notes(sl, "UDP ve TCP protokolleri baskın. Korelasyon analizinde ct_state_ttl, sttl ve rate özelliklerinin hedef değişkenle güçlü ilişki içinde olduğunu görüyoruz — bu özellikler modelimizde de öne çıkacak.")

# ── SLAYT 7: METODOLOJİ ─────────────────────────────────────
print("Slayt 7: Metodoloji")
sl = blank_slide(prs)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = C_WHITE
add_header(sl, "Metodoloji: Analiz Boru Hattı")
add_slide_number(sl, 7)

steps = [
    ("1  VERİ YÜKLEME",    "UNSW_NB15_training-set.csv  +  UNSW_NB15_testing-set.csv\n82,332 eğitim | 175,341 test | 42 özellik",   C_DARK_BLUE),
    ("2  ÖN İŞLEME",       "Label Encoding (proto, service, state)  •  StandardScaler\nEksik değer → medyan ile doldurma",            C_MED_BLUE),
    ("3  MODEL EĞİTİMİ",   "Random Forest: 100 ağaç, max_depth=20\nSınıf ağırlıklandırma: class_weight='balanced'",                  C_ACCENT_BLUE),
    ("4  DEĞERLENDİRME",   "Accuracy, F1-Score, ROC-AUC\nConfusion Matrix  •  Feature Importance",                                    RGBColor(0x1A, 0x6D, 0x37)),
    ("5  ETİK ANALİZ",     "FP/FN risk değerlendirmesi\nGDPR / Mahremiyet ve veri minimizasyonu incelemesi",                          RGBColor(0x7B, 0x1F, 0xA2)),
]
step_h = 0.88
step_w = SW - 1.0
label_w = 2.8
for i, (num, body, clr) in enumerate(steps):
    sy = HEADER_H + 0.1 + i * (step_h + 0.06)
    add_rect(sl, 0.5, sy, label_w, step_h, clr)
    add_text_box(sl, num,
                 0.55, sy + 0.12, label_w - 0.1, step_h - 0.24,
                 font_size=13, bold=True, color=C_WHITE,
                 align=PP_ALIGN.LEFT, margin_inches=0)
    add_rect(sl, 0.5 + label_w, sy, step_w - label_w, step_h, C_LIGHT_BLUE, clr, 1)
    add_text_box(sl, body,
                 0.5 + label_w + 0.15, sy + 0.1, step_w - label_w - 0.25, step_h - 0.2,
                 font_size=12, color=C_DARK_TEXT,
                 align=PP_ALIGN.LEFT, margin_inches=0)
    if i < len(steps) - 1:
        add_text_box(sl, "▼",
                     0.5 + label_w * 0.5 - 0.15, sy + step_h,
                     0.3, 0.1, font_size=9, color=clr,
                     align=PP_ALIGN.CENTER, margin_inches=0)
set_notes(sl, "Metodolojimiz beş adımdan oluşuyor. Veriyi yükledikten sonra kategorik değişkenleri sayısala çevirdik, standartlaştırdık ve eksik değerleri medyanla doldurduk. Model olarak Random Forest seçtik çünkü hem yüksek performanslı hem de yorumlanabilir bir modeldir. Son adımda salt teknik değerlendirmenin ötesine geçerek etik ve mahremiyet boyutunu analiz ettik.")

# ── SLAYT 8: PERFORMANS METRİKLERİ ─────────────────────────
print("Slayt 8: Performans Metrikleri")
sl = blank_slide(prs)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = C_WHITE
add_header(sl, "Model Sonuçları: Performans Metrikleri")
add_slide_number(sl, 8)

# 4 metrik kutu (2x2)
mw = (SW - 1.2) / 4 - 0.12
metrics = [
    ("DOĞRULUK",           "%90.13",  C_ACCENT_BLUE),
    ("F1 SKORU",           "0.9225",  RGBColor(0x1B, 0x6E, 0x35)),
    ("ROC-AUC",            "0.9864",  C_PURPLE),
    ("PRECISION (Saldırı)","% 99.0",  RGBColor(0xBF, 0x5C, 0x00)),
]
for i, (lbl, val, clr) in enumerate(metrics):
    mx = 0.5 + i * (mw + 0.18)
    metric_card(sl, mx, HEADER_H + 0.15, mw, 1.55, lbl, val, clr)

# Sınıflandırma raporu tablosu
add_text_box(sl, "Sınıflandırma Raporu (Test Seti: 175,341 kayıt)",
             0.5, HEADER_H + 1.88, SW - 1.0, 0.38,
             font_size=14, bold=True, color=C_DARK_BLUE, margin_inches=0)
headers = ["Sınıf", "Precision", "Recall", "F1-Score", "Support"]
rows = [
    ["Normal",        "0.77",   "0.98",   "0.86",   "56,000"],
    ["Saldırı",       "0.99",   "0.86",   "0.92",   "119,341"],
    ["Macro Ort.",    "0.88",   "0.92",   "0.89",   "175,341"],
    ["Ağırlıklı Ort.","0.92",   "0.90",   "0.90",   "175,341"],
]
add_table_styled(sl, headers, rows, 0.5, HEADER_H + 2.28, SW - 1.0, 2.95)
set_notes(sl, "Modelimiz test setinde yüzde doksan doğruluk elde etti. ROC-AUC değeri 0.986, bu ayrıştırma gücünün çok yüksek olduğunu gösteriyor. Saldırı sınıfında precision yüzde 99 — yani alarm verdiğinde neredeyse kesinlikle doğru. Ancak recall yüzde 86, yani saldırıların yüzde 14'ünü kaçırıyoruz. Bu güvenlik riski açısından kritik bir bulgudur.")

# ── SLAYT 9: CONFUSION MATRIX ───────────────────────────────
print("Slayt 9: Confusion Matrix")
sl = blank_slide(prs)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = C_WHITE
add_header(sl, "Karmaşıklık Matrisi (Confusion Matrix)")
add_slide_number(sl, 9)

add_image(sl, "confusion_matrix.png",
          0.3, HEADER_H + 0.1, SW - 0.6, 4.05)

# 4 sonuç kutusu
box_w = (SW - 1.0) / 4 - 0.1
boxes = [
    ("TP: 103,059",       "Saldırı doğru tespit",        RGBColor(0x1B, 0x6E, 0x35), C_GREEN_LIGHT),
    ("FN: 16,282 ⚠",     "Saldırı KAÇIRILDI",           C_RED,                       C_RED_LIGHT),
    ("FP: 1,023 ⚠",      "Normal → YANLIŞ işaretlendi", C_ORANGE,                    C_ORANGE_LIGHT),
    ("TN: 54,977",        "Normal doğru tespit",          C_ACCENT_BLUE,               C_LIGHT_BLUE),
]
for i, (val, lbl, txt_clr, bg_clr) in enumerate(boxes):
    bx = 0.5 + i * (box_w + 0.13)
    add_rect(sl, bx, SH - 1.18, box_w, 0.9, bg_clr, txt_clr, 1)
    add_text_box(sl, val,
                 bx + 0.05, SH - 1.15, box_w - 0.1, 0.38,
                 font_size=12, bold=True, color=txt_clr,
                 align=PP_ALIGN.CENTER, margin_inches=0)
    add_text_box(sl, lbl,
                 bx + 0.05, SH - 0.80, box_w - 0.1, 0.55,
                 font_size=10, color=txt_clr,
                 align=PP_ALIGN.CENTER, margin_inches=0)
set_notes(sl, "Confusion matrix bize iki önemli hata türünü gösteriyor. 16 bin 282 yanlış negatif — gerçek saldırı atlandı, bu bir güvenlik riskidir. 1023 yanlış pozitif — normal bir kullanıcı saldırgan olarak işaretlendi, bu etik ve operasyonel bir sorundur. Bu iki hata tipinin farklı risk kategorileri taşıdığını vurgulamak istiyorum.")

# ── SLAYT 10: ROC & FEATURE IMPORTANCE ──────────────────────
print("Slayt 10: ROC & Feature Importance")
sl = blank_slide(prs)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = C_WHITE
add_header(sl, "ROC Eğrisi ve Özellik Önemi")
add_slide_number(sl, 10)

half = (SW - 0.9) / 2
add_image(sl, "roc_curve.png",
          0.3, HEADER_H + 0.1, half, 5.5)
add_image(sl, "feature_importance.png",
          0.3 + half + 0.3, HEADER_H + 0.1, half, 5.5)

add_text_box(sl, "AUC = 0.9864",
             0.5, SH - 0.42, half - 0.3, 0.35,
             font_size=12, bold=True, color=C_ACCENT_BLUE,
             align=PP_ALIGN.CENTER, margin_inches=0)
add_text_box(sl, "ct_state_ttl, sttl, dttl en belirleyici özellikler",
             0.5 + half + 0.3, SH - 0.42, half - 0.3, 0.35,
             font_size=12, bold=True, color=C_ACCENT_BLUE,
             align=PP_ALIGN.CENTER, margin_inches=0)
set_notes(sl, "ROC eğrisi neredeyse mükemmel — AUC 0.986 modelimizin rastgele tahminden çok üstün olduğunu kanıtlıyor. Özellik önemi grafiğinde ise ct_state_ttl, sttl ve dttl özelliklerinin en belirleyici olduğunu görüyoruz. Bu özellikler ağ bağlantısının durumu ve yaşam süresi ile ilgili — saldırı trafiği bu açıdan normal trafikten belirgin şekilde ayrışıyor.")

# ── SLAYT 11: ETİK RİSK ANALİZİ ────────────────────────────
print("Slayt 11: Etik Risk Analizi")
sl = blank_slide(prs)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = C_WHITE
add_header(sl, "Etik ve Güvenlik Risk Analizi")
add_slide_number(sl, 11)

add_image(sl, "risk_analysis.png",
          0.3, HEADER_H + 0.1, SW - 0.6, 3.55)

# İki açıklama kutusu yan yana
half_w = (SW - 1.1) / 2
# Etik risk (turuncu)
add_rect(sl, 0.4, HEADER_H + 3.8, half_w, 2.8, C_ORANGE_LIGHT, C_ORANGE, 2)
add_text_box(sl, "ETİK RİSK",
             0.55, HEADER_H + 3.85, half_w - 0.3, 0.38,
             font_size=14, bold=True, color=C_ORANGE, margin_inches=0)
eth_bullets = [
    "1,023 masum kullanıcı saldırgan olarak işaretlendi",
    "Haksız erişim engelleme riski",
    "Kullanıcıda güven kaybı",
    "Orantısız izleme sorunu (GDPR)",
]
add_bullet_box(sl, eth_bullets, 0.55, HEADER_H + 4.28, half_w - 0.3, 2.2, font_size=12, color=C_ORANGE)

# Güvenlik riski (kırmızı)
add_rect(sl, 0.5 + half_w + 0.1, HEADER_H + 3.8, half_w, 2.8, C_RED_LIGHT, C_RED, 2)
add_text_box(sl, "GÜVENLİK RİSKİ",
             0.65 + half_w + 0.1, HEADER_H + 3.85, half_w - 0.3, 0.38,
             font_size=14, bold=True, color=C_RED, margin_inches=0)
sec_bullets = [
    "16,282 gerçek saldırı tespit edilemedi",
    "Sistem güvenlik açığı oluşabilir",
    "Veri ihlali ve ağ ihlali riski",
    "Altyapı güvenliği tehlikeye girer",
]
add_bullet_box(sl, sec_bullets, 0.65 + half_w + 0.1, HEADER_H + 4.28, half_w - 0.3, 2.2, font_size=12, color=C_RED)
set_notes(sl, "Bu grafik dört kategoriyi karşılaştırıyor. Turuncu sütun — yanlış pozitifler — etik riski temsil ediyor: masum bir kullanıcı haksız yere tehdit olarak işaretleniyor. Kırmızı sütun — yanlış negatifler — güvenlik riskini temsil ediyor: gerçek bir saldırı gözden kaçıyor. İdeal bir sistemde bu iki hata türü arasında dengeli bir trade-off kurulmalıdır.")

# ── SLAYT 12: MAHREMİYET & GDPR ─────────────────────────────
print("Slayt 12: Mahremiyet & GDPR")
sl = blank_slide(prs)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = C_WHITE
add_header(sl, "Mahremiyet ve GDPR Değerlendirmesi")
add_slide_number(sl, 12)

headers = ["Özellik / Veri",   "Açıklama",               "GDPR Riski"]
rows = [
    ["IP Adresi",          "Kaynak / Hedef adresi",          "YÜKSEK — Kişisel veri"],
    ["Port Numarası",      "Kaynak / Hedef port",            "ORTA — Dolaylı tanımlayıcı"],
    ["Zaman Damgası",      "Bağlantı başlangıç / bitiş",     "ORTA — Davranış profili"],
    ["Protokol",           "TCP / UDP / HTTP vb.",           "DÜŞÜK"],
    ["ct_src_ltm",         "Son bağlantı sayısı (kaynak)",   "ORTA — Kullanım örüntüsü"],
    ["ct_dst_ltm",         "Son bağlantı sayısı (hedef)",    "ORTA — Kullanım örüntüsü"],
]
add_table_styled(sl, headers, rows, 0.4, HEADER_H + 0.15, SW - 0.8, 3.5)

# GDPR ilkeleri — 3 kutu
prin_w = (SW - 1.2) / 3 - 0.1
principles = [
    ("Amaçla Sınırlılık",    "Veriler yalnızca güvenlik amacıyla işlenmeli, başka amaçlarla kullanılmamalı",         C_DARK_BLUE),
    ("Veri Minimizasyonu",    "Saldırı tespiti için zorunlu olmayan kişisel veriler toplanmamalı ve işlenmemeli",      C_ACCENT_BLUE),
    ("Orantılılık",           "İzleme kapsamı ve yoğunluğu, mevcut tehditle orantılı olmalı; aşırı izlemeden kaçınılmalı", C_MED_BLUE),
]
for i, (ptitle, pbody, pclr) in enumerate(principles):
    px = 0.4 + i * (prin_w + 0.2)
    add_rect(sl, px, HEADER_H + 3.85, prin_w, 0.5, pclr)
    add_text_box(sl, f"GDPR: {ptitle}",
                 px + 0.08, HEADER_H + 3.88, prin_w - 0.16, 0.42,
                 font_size=12, bold=True, color=C_WHITE,
                 align=PP_ALIGN.LEFT, margin_inches=0)
    add_rect(sl, px, HEADER_H + 4.38, prin_w, 1.7, C_LIGHT_BLUE, pclr, 1)
    add_text_box(sl, pbody,
                 px + 0.1, HEADER_H + 4.45, prin_w - 0.2, 1.55,
                 font_size=11, color=C_DARK_TEXT,
                 align=PP_ALIGN.LEFT, margin_inches=0)
set_notes(sl, "Saldırı tespiti sistemleri IP adresleri, zaman damgaları ve bağlantı örüntüleri gibi kişisel veri niteliği taşıyabilecek bilgileri işler. GDPR bu tür verilerin işlenmesinde amaçla sınırlılık, veri minimizasyonu ve orantılılık ilkelerine uyulmasını zorunlu kılmaktadır. Sistemin tasarımında bu ilkeler gözetilmezse ciddi uyum ve etik sorunları ortaya çıkabilir.")

# ── SLAYT 13: ETİK ÇERÇEVE - TRADE-OFF ──────────────────────
print("Slayt 13: Etik Çerçeve")
sl = blank_slide(prs)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = C_WHITE
add_header(sl, "Etik Çerçeve: İki Taraflı Risk Dengesi")
add_slide_number(sl, 13)

# Sol kutu
add_rect(sl, 0.4, HEADER_H + 0.2, 4.4, 3.2, C_RED_LIGHT, C_RED, 2)
add_text_box(sl, "RECALL'U ARTIR\n(Daha az FN)",
             0.55, HEADER_H + 0.3, 4.1, 0.7,
             font_size=16, bold=True, color=C_RED,
             align=PP_ALIGN.CENTER, margin_inches=0)
recall_bullets = [
    "Daha az saldırı kaçırılır",
    "Daha fazla yanlış alarm üretilir",
    "Masum kullanıcılar etkilenir",
    "Mahremiyet / adalet riski artar",
]
add_bullet_box(sl, recall_bullets, 0.6, HEADER_H + 1.1, 4.0, 2.2, font_size=12, color=C_RED)

# Orta - Trade-off
add_rect(sl, 5.07, HEADER_H + 0.55, 3.2, 2.7, C_DARK_BLUE)
add_text_box(sl, "TRADE-OFF\n\nEşik değeri\n(threshold)\nayarlanarak\ndenge kurulur",
             5.12, HEADER_H + 0.65, 3.1, 2.5,
             font_size=14, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER, margin_inches=0)

# Sağ kutu
add_rect(sl, 8.53, HEADER_H + 0.2, 4.4, 3.2, C_ORANGE_LIGHT, C_ORANGE, 2)
add_text_box(sl, "PRECISION'I ARTIR\n(Daha az FP)",
             8.68, HEADER_H + 0.3, 4.1, 0.7,
             font_size=16, bold=True, color=C_ORANGE,
             align=PP_ALIGN.CENTER, margin_inches=0)
prec_bullets = [
    "Daha az yanlış alarm üretilir",
    "Daha fazla saldırı kaçırılır",
    "Güvenlik açığı riski artar",
    "Ağ ihlali tehlikesi büyür",
]
add_bullet_box(sl, prec_bullets, 8.68, HEADER_H + 1.1, 4.1, 2.2, font_size=12, color=C_ORANGE)

# Alttaki açıklama kutuları
add_rect(sl, 0.4, HEADER_H + 3.55, (SW - 1.0) / 2 - 0.1, 2.7, C_LIGHT_BLUE, C_ACCENT_BLUE, 1)
add_text_box(sl, "Teknik Boyut",
             0.55, HEADER_H + 3.6, (SW - 1.2) / 2 - 0.1, 0.38,
             font_size=13, bold=True, color=C_ACCENT_BLUE, margin_inches=0)
add_text_box(sl, "Eşik değeri (threshold) ayarlanarak precision-recall dengesi değiştirilebilir. Yüksek güvenlik ortamlarında recall önceliklendirilmeli; kullanıcı gizliliğinin ön planda olduğu ortamlarda ise precision artırılmalıdır.",
             0.55, HEADER_H + 4.02, (SW - 1.2) / 2 - 0.1, 2.1,
             font_size=12, color=C_DARK_TEXT, margin_inches=0)

add_rect(sl, 0.5 + (SW - 1.0) / 2, HEADER_H + 3.55, (SW - 1.0) / 2 - 0.1, 2.7, C_PURPLE_LIGHT, C_PURPLE, 1)
add_text_box(sl, "Etik Boyut",
             0.65 + (SW - 1.0) / 2, HEADER_H + 3.6, (SW - 1.2) / 2 - 0.1, 0.38,
             font_size=13, bold=True, color=C_PURPLE, margin_inches=0)
add_text_box(sl, "Otomatik karar sistemleri şeffaflık, hesap verebilirlik ve açıklanabilirlik (XAI) gerektirir. Her yanlış sınıflandırma, somut bir insanı etkiler; bu nedenle model kararları açıklanabilir olmalıdır.",
             0.65 + (SW - 1.0) / 2, HEADER_H + 4.02, (SW - 1.2) / 2 - 0.1, 2.1,
             font_size=12, color=C_DARK_TEXT, margin_inches=0)
set_notes(sl, "Burada temel bir etik ikilemle karşı karşıyayız. Recall'u artırmak — yani daha az saldırı kaçırmak — kaçınılmaz olarak daha fazla yanlış alarmı beraberinde getirir ve masum kullanıcıları etkiler. Tersine precision'ı artırmak daha az yanlış alarm anlamına gelir ama bazı gerçek saldırılar atlanır. Bu denge, sistemin kullanıldığı bağlama ve risk iştahına göre belirlenmelidir.")

# ── SLAYT 14: SONUÇ ─────────────────────────────────────────
print("Slayt 14: Sonuç ve Öneriler")
sl = blank_slide(prs)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = C_DARK_BLUE

# Başlık
add_rect(sl, 0, 0, SW, 1.05, RGBColor(0x0D, 0x47, 0xA1))
add_text_box(sl, "Sonuç ve Öneriler",
             0.4, 0.12, SW - 0.8, 0.82,
             font_size=28, bold=True, color=C_WHITE,
             align=PP_ALIGN.LEFT, margin_inches=0)

half_w = (SW - 1.2) / 2

# Bulgular kutusu
add_rect(sl, 0.4, 1.2, half_w, 5.4, RGBColor(0x0D, 0x47, 0xA1))
add_text_box(sl, "BULGULAR",
             0.55, 1.28, half_w - 0.3, 0.45,
             font_size=15, bold=True, color=RGBColor(0x90, 0xCA, 0xF9),
             align=PP_ALIGN.LEFT, margin_inches=0)
findings = [
    ("✓", "Random Forest %90.13 doğrulukla güçlü performans sergiledi",        C_GREEN_LIGHT),
    ("✓", "ROC-AUC: 0.9864 — yüksek ayrıştırma gücü elde edildi",             C_GREEN_LIGHT),
    ("✓", "103,059 saldırı başarıyla tespit edildi",                             C_GREEN_LIGHT),
    ("⚠", "16,282 saldırı kaçırıldı  →  Güvenlik riski",                       RGBColor(0xFF, 0xCC, 0x80)),
    ("⚠", "1,023 kullanıcı hatalı etiketlendi  →  Etik risk",                  RGBColor(0xFF, 0xCC, 0x80)),
]
for i, (icon, text, clr) in enumerate(findings):
    fy = 1.78 + i * 0.85
    add_text_box(sl, icon,
                 0.55, fy, 0.35, 0.6,
                 font_size=18, bold=True, color=clr,
                 align=PP_ALIGN.CENTER, margin_inches=0)
    add_text_box(sl, text,
                 0.95, fy + 0.05, half_w - 0.85, 0.65,
                 font_size=13, color=C_WHITE,
                 align=PP_ALIGN.LEFT, margin_inches=0)

# Öneriler kutusu
add_rect(sl, 0.6 + half_w, 1.2, half_w, 5.4, RGBColor(0x1E, 0x2E, 0x8E))
add_text_box(sl, "ÖNERİLER",
             0.75 + half_w, 1.28, half_w - 0.3, 0.45,
             font_size=15, bold=True, color=RGBColor(0x90, 0xCA, 0xF9),
             align=PP_ALIGN.LEFT, margin_inches=0)
recs = [
    "Eşik değeri bağlama göre optimize edilmeli",
    "Veri minimizasyonu prensibi uygulanmalı",
    "IP adresleri anonimleştirilmeli",
    "Model kararları açıklanabilir olmalı (XAI)",
    "Periyodik denetim ve model güncellemesi yapılmalı",
    "Etik gözetim mekanizması kurulmalı",
]
for i, r in enumerate(recs):
    ry = 1.78 + i * 0.82
    add_text_box(sl, f"→  {r}",
                 0.75 + half_w, ry, half_w - 0.5, 0.72,
                 font_size=13, color=C_WHITE,
                 align=PP_ALIGN.LEFT, margin_inches=0)

# Alt banner
add_rect(sl, 0, SH - 0.9, SW, 0.9, RGBColor(0x0A, 0x37, 0x6E))
add_text_box(sl, '"Büyük veri güvenlik sistemleri; teknik başarının yanı sıra etik, adalet ve mahremiyet ilkeleriyle de uyumlu tasarlanmalıdır."',
             0.4, SH - 0.85, SW - 0.8, 0.78,
             font_size=14, bold=True, italic=True, color=RGBColor(0xBB, 0xDE, 0xFB),
             align=PP_ALIGN.CENTER, margin_inches=0)

add_slide_number(sl, 14)
set_notes(sl, "Özetlemek gerekirse: Random Forest modeli teknik açıdan güçlü sonuçlar verdi. Ancak yanlış sınıflandırmalar hem güvenlik hem etik açıdan gerçek riskler doğuruyor. En önemli çıkarımım şu: Büyük veri güvenlik sistemleri sadece doğruluk oranıyla değil, verdiği kararların insan üzerindeki etkisiyle de değerlendirilmelidir. Teşekkür ederim, sorularınızı alabilirim.")

# ─── KAYDET ──────────────────────────────────────────────────
prs.save(OUT)
print(f"\n✅  Sunum kaydedildi: {OUT}")
print(f"   Toplam slayt: {len(prs.slides)}")
