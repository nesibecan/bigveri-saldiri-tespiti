"""
Slayt 7 ve Slayt 13'ü hedefli olarak güncelle:
  - Slayt 7: MODEL EĞİTİMİ adımına "Neden RF?" gerekçesi ekle
  - Slayt 13: Threshold/eşik değeri kararını öne çıkaran callout kutusu ekle
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = r"C:\Users\seyma\OneDrive\Masaüstü\bigdataproje"
PATH = os.path.join(BASE, "BigVeri_SaldiriTespiti_Sunum.pptx")

prs = Presentation(PATH)
SW, SH = 13.333, 7.5
HEADER_H = 0.85

C_DARK_BLUE   = RGBColor(0x1A, 0x23, 0x7E)
C_ACCENT_BLUE = RGBColor(0x15, 0x65, 0xC0)
C_MED_BLUE    = RGBColor(0x28, 0x35, 0x93)
C_LIGHT_BLUE  = RGBColor(0xE3, 0xF2, 0xFD)
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK_TEXT   = RGBColor(0x1A, 0x23, 0x7E)
C_PURPLE      = RGBColor(0x4A, 0x14, 0x8C)
C_PURPLE_LIGHT= RGBColor(0xF3, 0xE5, 0xF5)
C_GOLD        = RGBColor(0xFF, 0xC1, 0x07)
C_GOLD_DARK   = RGBColor(0xF5, 0x7F, 0x17)
C_GOLD_LIGHT  = RGBColor(0xFF, 0xF9, 0xC4)

def add_rect(slide, x, y, w, h, fill_color, line_color=None, lw=0):
    shape = slide.shapes.add_shape(1,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(lw)
    else:
        shape.line.fill.background()
    return shape

def add_tb(slide, text, x, y, w, h,
           size=13, bold=False, italic=False,
           color=None, align=PP_ALIGN.LEFT, face="Calibri"):
    if color is None:
        color = C_DARK_TEXT
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top  = tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name   = face
    return tb

# ════════════════════════════════════════
# SLAYT 7 — Metodoloji (index 6)
# Adım 3 bloğunun yanına "Neden RF?" bilgi kutusu ekle
# ════════════════════════════════════════
sl7 = prs.slides[6]
print("Slayt 7 güncelleniyor...")

# Sağ kenara ayrı bir "Neden Random Forest?" bilgi paneli
# (mevcut adım kutuları x=0.5..3.3 bölgesinde; sağ alan: 3.5 .. 12.8)
panel_x = 3.7
panel_y = HEADER_H + 0.12
panel_w = SW - panel_x - 0.4
panel_h = 6.1

add_rect(sl7, panel_x, panel_y, panel_w, panel_h,
         RGBColor(0xE8, 0xEA, 0xF6), C_DARK_BLUE, 1)

add_tb(sl7, "Neden Random Forest?",
       panel_x + 0.12, panel_y + 0.12, panel_w - 0.24, 0.45,
       size=15, bold=True, color=C_DARK_BLUE)

reasons = [
    ("Aşırı öğrenmeye direnç",
     "Yüzlerce ağacın oylama ortalaması alındığından tek ağaca göre çok daha kararlıdır."),
    ("Dağılım varsayımı yok",
     "Lojistik regresyon gibi özellik dağılımı varsayımına ihtiyaç duymaz."),
    ("Karma özellik tipleri",
     "Sayısal + kodlanmış kategorik özellikleri aynı anda işleyebilir."),
    ("Yorumlanabilirlik",
     "Feature importance skoru sayesinde hangi özelliğin ne kadar katkı verdiği ölçülür."),
    ("StandardScaler notu",
     "RF ölçeğe duyarsızdır; scaler pipeline standardizasyonu\niçin eklenmiştir, modelin zorunluluğundan değil."),
]

ry = panel_y + 0.65
for title, body in reasons:
    add_rect(sl7, panel_x + 0.12, ry, panel_w - 0.24, 0.08,
             C_ACCENT_BLUE)
    add_tb(sl7, title,
           panel_x + 0.12, ry + 0.06, panel_w - 0.24, 0.30,
           size=12, bold=True, color=C_ACCENT_BLUE)
    add_tb(sl7, body,
           panel_x + 0.20, ry + 0.35, panel_w - 0.40, 0.50,
           size=11, color=C_DARK_TEXT)
    ry += 0.98

print("  -> Slayt 7 guncellendi.")

# ════════════════════════════════════════
# SLAYT 13 — Etik Trade-off (index 12)
# Alt yarıya öne çıkan threshold callout kutusu ekle
# ════════════════════════════════════════
sl13 = prs.slides[12]
print("Slayt 13 güncelleniyor...")

# Mevcut iki alttaki kutu y ≈ 4.4 civarında; onların altına sarı callout
callout_y = SH - 1.10
callout_h = 0.85

add_rect(sl13, 0.4, callout_y - 0.05, SW - 0.8, callout_h + 0.1,
         C_GOLD_DARK)
add_rect(sl13, 0.42, callout_y - 0.03, SW - 0.84, callout_h + 0.06,
         C_GOLD_LIGHT, C_GOLD_DARK, 2)

add_tb(sl13, "TEMEL ÇIKARIM:",
       0.6, callout_y + 0.04, 2.0, 0.38,
       size=12, bold=True, color=C_GOLD_DARK)

add_tb(sl13,
       "Eşik değeri (varsayılan 0.5) sabit tutulmamalıdır. "
       "Güvenlik öncelikli ortamlarda eşik düşürülerek recall artırılır (daha az kaçan saldırı); "
       "mahremiyet öncelikli ortamlarda eşik yükseltilerek yanlış pozitif azaltılır. "
       "Bu seçim bağlama özel, bilinçli ve etik olarak gerekçelendirilmiş bir karar olmalıdır.",
       2.7, callout_y + 0.04, SW - 3.1, 0.75,
       size=11, bold=False, italic=True, color=C_DARK_TEXT)

print("  -> Slayt 13 guncellendi.")

prs.save(PATH)
print(f"\nDosya kaydedildi: {PATH}")
